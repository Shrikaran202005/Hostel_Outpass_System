import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Colors
    COLOR_BG = RGBColor(15, 23, 42)        # Slate 900
    COLOR_CARD = RGBColor(30, 41, 59)      # Slate 800
    COLOR_TEXT = RGBColor(241, 245, 249)   # Slate 100
    COLOR_MUTED = RGBColor(148, 163, 184)  # Slate 400
    COLOR_ACCENT = RGBColor(124, 58, 237)  # Purple 600
    COLOR_EMERALD = RGBColor(16, 185, 129) # Emerald 500
    COLOR_ROSE = RGBColor(244, 63, 94)     # Rose 500

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="HOSTEL OUTING PERMISSION SYSTEM"):
        # Header category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT

    def add_card(slide, left, top, width, height, title, items, border_color=COLOR_CARD):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)

        tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(width - 0.4), Inches(height - 0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_TEXT
        p0.space_after = Pt(12)

        for item in items:
            p = tf.add_paragraph()
            p.text = f"•  {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = COLOR_MUTED
            p.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)

    t_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
    tf1 = t_box.text_frame
    tf1.word_wrap = True

    p_sub = tf1.paragraphs[0]
    p_sub.text = "COLLEGIATE CAMPUS AUTOMATION"
    p_sub.font.size = Pt(14)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_ACCENT
    p_sub.space_after = Pt(10)

    p_title = tf1.add_paragraph()
    p_title.text = "Hostel Outing Permission &\nApproval Management System"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT
    p_title.space_after = Pt(16)

    p_desc = tf1.add_paragraph()
    p_desc.text = "Stage 3 AI-Assisted Development Loop & Workflow Automation"
    p_desc.font.size = Pt(18)
    p_desc.font.color.rgb = COLOR_MUTED

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_header(slide2, "The Real-World Hostel Outing Problem")

    add_card(slide2, 0.8, 1.8, 3.6, 5.0, "Manual Paperwork", [
        "Handwritten physical outing slips.",
        "High risk of lost or forged permission forms.",
        "No centralized audit database."
    ])
    add_card(slide2, 4.8, 1.8, 3.6, 5.0, "Coordination Delays", [
        "Students wait hours outside HOD/Warden offices.",
        "Lack of parent verification tracking.",
        "Unclear status updates for students."
    ])
    add_card(slide2, 8.8, 1.8, 3.6, 5.0, "Security & Gate Gaps", [
        "Watchman cannot easily verify sign-offs.",
        "Late returns are difficult to log accurately.",
        "No role-scoped historical reporting."
    ])

    # -------------------------------------------------------------
    # SLIDE 3: Solution
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_header(slide3, "End-to-End Automated Solution Workflow")

    add_card(slide3, 0.8, 1.8, 2.7, 5.0, "1. Student Request", [
        "Online submission.",
        "Auto department & block mapping.",
        "Time window validation."
    ])
    add_card(slide3, 3.8, 1.8, 2.7, 5.0, "2. HOD Review", [
        "Academic approval.",
        "Strict department scope.",
        "Approve or reject."
    ])
    add_card(slide3, 6.8, 1.8, 2.7, 5.0, "3. Warden Consent", [
        "Mandatory parent phone check.",
        "Hostel block scope.",
        "Final approval grant."
    ])
    add_card(slide3, 9.8, 1.8, 2.7, 5.0, "4. Gate Verification", [
        "Watchman verification.",
        "Record Exit & Return.",
        "Late Return Detection."
    ])

    # -------------------------------------------------------------
    # SLIDE 4: Architecture
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4)
    add_header(slide4, "Production Tech Stack & Architecture")

    add_card(slide4, 0.8, 1.8, 3.6, 5.0, "Frontend Layer", [
        "React + Vite + TypeScript.",
        "Vanilla CSS / Tailwind styling.",
        "Role-based routing & protected views."
    ])
    add_card(slide4, 4.8, 1.8, 3.6, 5.0, "Backend Layer", [
        "FastAPI framework.",
        "OAuth2 JWT Bearer authentication.",
        "SQLAlchemy ORM + Pydantic validation."
    ])
    add_card(slide4, 8.8, 1.8, 3.6, 5.0, "Data & Security", [
        "Relational SQL schema.",
        "Unified audit & gate logs.",
        "Strict role-based authorization."
    ])

    # -------------------------------------------------------------
    # SLIDE 5: Role-Based Access
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5)
    add_header(slide5, "Strict Role-Based Authorization Scoping")

    add_card(slide5, 0.8, 1.8, 2.7, 5.0, "STUDENT", [
        "Submit outing requests.",
        "Cancel pending requests.",
        "View own history only."
    ])
    add_card(slide5, 3.8, 1.8, 2.7, 5.0, "HOD", [
        "Academic approval.",
        "Strict Department Scope.",
        "View department history."
    ])
    add_card(slide5, 6.8, 1.8, 2.7, 5.0, "WARDEN", [
        "Parent phone consent check.",
        "Strict Hostel Block Scope.",
        "View block history."
    ])
    add_card(slide5, 9.8, 1.8, 2.7, 5.0, "WATCHMAN", [
        "Gate security desk.",
        "Search all campus students.",
        "Record Exit & Return."
    ])

    # -------------------------------------------------------------
    # SLIDE 6: Stage 3 AI Loop
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6)
    add_header(slide6, "Stage 3 AI Development Loop (Late Return Detection)")

    add_card(slide6, 0.8, 1.8, 3.6, 5.0, "1. New Requirement", [
        "Feature: Late Return Detection.",
        "Actual Return > Expected Return.",
        "Set status = LATE_RETURN."
    ], border_color=COLOR_ACCENT)

    add_card(slide6, 4.8, 1.8, 3.6, 5.0, "2. Deliberate RED Run", [
        "New unit test created.",
        "Pytest output: 1 FAILED.",
        "Expected LATE_RETURN, got COMPLETED."
    ], border_color=COLOR_ROSE)

    add_card(slide6, 8.8, 1.8, 3.6, 5.0, "3. AI Debug & GREEN Fix", [
        "AI analyzed root cause.",
        "Updated watchman_record_return.",
        "Pytest & E2E: 100% PASSED."
    ], border_color=COLOR_EMERALD)

    # -------------------------------------------------------------
    # SLIDE 7: Empirical Testing Metrics
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_bg(slide7)
    add_header(slide7, "Empirical Verification & Test Results")

    add_card(slide7, 0.8, 1.8, 3.6, 5.0, "Backend Unit Tests", [
        "74 Passed / 0 Failed.",
        "Pytest execution: 28.54s.",
        "Role, workflow & security tests."
    ])
    add_card(slide7, 4.8, 1.8, 3.6, 5.0, "Playwright E2E Tests", [
        "8 Passed / 0 Failed.",
        "E2E execution: 26.2s.",
        "Full workflow & directory tests."
    ])
    add_card(slide7, 8.8, 1.8, 3.6, 5.0, "Production Build", [
        "TypeScript tsc: 0 Errors.",
        "Vite build: SUCCESS in 2.98s.",
        "Zero regressions."
    ])

    # -------------------------------------------------------------
    # SLIDE 8: Live Demo Plan
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_bg(slide8)
    add_header(slide8, "5-Minute Demonstration Plan")

    add_card(slide8, 0.8, 1.8, 5.6, 5.0, "2-Min Architecture Overview", [
        "0:00 - 0:30: Problem & manual permission delays.",
        "0:30 - 1:00: Multi-role routing & approval flow.",
        "1:00 - 1:30: Role scoping (HOD Dept / Warden Block).",
        "1:30 - 2:00: Stage 3 AI Loop & Late Return rules."
    ])
    add_card(slide8, 6.8, 1.8, 5.6, 5.0, "3-Min Live System Demo", [
        "2:00 - 2:30: Student creates request.",
        "2:30 - 2:50: HOD grants academic approval.",
        "2:50 - 3:10: Warden confirms parent & final approves.",
        "3:10 - 4:30: Watchman records exit & late return.",
        "4:30 - 5:00: Show LATE_RETURN in HOD/Warden History."
    ])

    # -------------------------------------------------------------
    # SLIDE 9: Conclusion
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    add_bg(slide9)
    add_header(slide9, "Conclusion & Project Accomplishments")

    add_card(slide9, 0.8, 1.8, 3.6, 5.0, "Automation Delivered", [
        "Complete automated approval workflow.",
        "Parent confirmation enforcement.",
        "Gate security desk integration."
    ])
    add_card(slide9, 4.8, 1.8, 3.6, 5.0, "Security & Scoping", [
        "Strict department & block isolation.",
        "Backend JWT role enforcement.",
        "Unified chronological audit timeline."
    ])
    add_card(slide9, 8.8, 1.8, 3.6, 5.0, "Stage 3 AI Success", [
        "Late Return Detection integrated.",
        "Verified via RED -> FIX -> GREEN loop.",
        "100% automated test coverage."
    ])

    os.makedirs("c:/Data/Inter Assign/presentation", exist_ok=True)
    out_path = "c:/Data/Inter Assign/presentation/Tactive_Assessment.pptx"
    prs.save(out_path)
    print(f"Presentation created successfully at {out_path}")

if __name__ == "__main__":
    create_presentation()
